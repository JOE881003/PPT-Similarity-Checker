import streamlit as st
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# ---------- 函式 ----------
def extract_text_from_ppt(uploaded_file) -> str:
    """從上傳的 PPT 檔案提取文字"""
    prs = Presentation(uploaded_file)
    text_list = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_list.append(shape.text)
    return " ".join(text_list)

def calculate_similarity(text_main: str, text_other: str) -> float:
    """計算兩份文字的相似度"""
    vectorizer = TfidfVectorizer().fit([text_main, text_other])
    vecs = vectorizer.transform([text_main, text_other])
    return cosine_similarity(vecs[0], vecs[1])[0][0]

# ---------- Streamlit UI ----------
st.title("PPT 相似度比對工具")

st.write("上傳兩份 **PPTX** 簡報，系統會幫你計算文字內容的相似度。")

col1, col2 = st.columns(2)

with col1:
    ppt1 = st.file_uploader("上傳第一份 PPT", type=["pptx"], key="ppt1")
with col2:
    ppt2 = st.file_uploader("上傳第二份 PPT", type=["pptx"], key="ppt2")

if ppt1 and ppt2:
    with st.spinner("正在分析簡報內容..."):
        text1 = extract_text_from_ppt(ppt1)
        text2 = extract_text_from_ppt(ppt2)
        score = calculate_similarity(text1, text2)

    st.success(f"兩份簡報的相似度為： **{score:.4f}**")
    
    # 顯示一些分析結果
    st.subheader("📌 文字摘要")
    st.write("**簡報 1 部分內容：**", text1[:300] + "...")
    st.write("**簡報 2 部分內容：**", text2[:300] + "...")
else:
    st.info("請先上傳兩份 PPTX 檔案")