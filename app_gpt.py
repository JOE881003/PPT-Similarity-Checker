import streamlit as st
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# ==== 新增：OpenAI ====
from openai import OpenAI

# ---------- 你原本的函式 ----------
def extract_text_from_ppt(uploaded_file) -> str:
    """從上傳的 PPT 檔案提取文字"""
    prs = Presentation(uploaded_file)
    text_list = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_list.append(shape.text)
    return " ".join(text_list)

def calculate_similarity(text_main: str, text_other: str) -> float:
    """計算兩份文字的相似度"""
    vectorizer = TfidfVectorizer().fit([text_main, text_other])
    vecs = vectorizer.transform([text_main, text_other])
    return cosine_similarity(vecs[0], vecs[1])[0][0]

# ---------- 新增：摘要工具 ----------
def chunk_text(s: str, max_chars: int = 4000):
    """粗略依字元長度分段（避免一次丟太長）"""
    s = s.replace("\x00", " ")  # 清理怪字元
    return [s[i:i+max_chars] for i in range(0, len(s), max_chars)] or ["(空白)"]

def summarize_long_text(text: str, client: OpenAI, model: str = "gpt-4o-mini") -> str:
    """先對每段做局部摘要，再做全局整合摘要"""
    parts = chunk_text(text, max_chars=4000)

    partial_summaries = []
    for idx, part in enumerate(parts, 1):
        prompt = (
            "你是一位簡報顧問，請用繁體中文將下列內容重點式摘要成 3-6 點，"
            "每點一句話、避免重複與冗詞，並在最後給 3-6 個關鍵字（#標籤）。\n\n"
            f"[分段 {idx}/{len(parts)} 內容]\n{part}"
        )
        resp = client.responses.create(model=model, input=prompt)
        partial_summaries.append(resp.output_text.strip())

    # 全局整合
    combine_prompt = (
        "以下是多段的局部摘要，請融合成一份『最終簡報摘要』：\n"
        "1) 先給一段 2-3 句的總覽（TL;DR）。\n"
        "2) 接著列 5-8 點條列重點（每點 ≤ 25 字）。\n"
        "3) 最後列出 5-10 個關鍵字（#標籤）。\n\n"
        "=== 局部摘要開始 ===\n" + "\n\n---\n\n".join(partial_summaries)
    )
    final_resp = client.responses.create(model=model, input=combine_prompt)
    return final_resp.output_text.strip()

# ---------- Streamlit UI ----------
st.title("PPT 相似度比對工具 + GPT 摘要")
st.write("上傳兩份 **PPTX** 簡報，系統會計算文字相似度，並可用 **GPT** 產生摘要。")

# 建議用 Secrets；這裡提供臨時輸入欄位以便測試
#api_key = st.text_input("OpenAI API Key（以 sk- 開頭）", type="password", help="正式上線請改用 st.secrets 保存")
api_key = st.secrets["openai"]["api_key"]
model = st.selectbox("摘要模型", ["gpt-4o-mini", "gpt-4o"], index=0)

col1, col2 = st.columns(2)
with col1:
    ppt1 = st.file_uploader("上傳第一份 PPT", type=["pptx"], key="ppt1")
with col2:
    ppt2 = st.file_uploader("上傳第二份 PPT", type=["pptx"], key="ppt2")

if ppt1 and ppt2:
    with st.spinner("正在分析簡報內容..."):
        text1 = extract_text_from_ppt(ppt1)
        text2 = extract_text_from_ppt(ppt2)
        score = calculate_similarity(text1, text2)

    st.success(f"兩份簡報的相似度為： **{score:.4f}**")

    st.subheader("文字摘要（截斷展示）")
    st.write("**簡報 1 部分內容：**", (text1[:500] + "...") if len(text1) > 500 else text1)
    st.write("**簡報 2 部分內容：**", (text2[:500] + "...") if len(text2) > 500 else text2)

    # ---------- 新增：用 GPT 產生摘要 ----------
    if api_key:
        if st.button("產生兩份 PPT 的 GPT 摘要"):
            try:
                client = OpenAI(api_key=api_key)
                with st.spinner("GPT 正在產生摘要（簡報 1）…"):
                    summary1 = summarize_long_text(text1, client, model=model)
                with st.spinner("GPT 正在產生摘要（簡報 2）…"):
                    summary2 = summarize_long_text(text2, client, model=model)

                st.subheader("GPT 摘要：簡報 1")
                st.markdown(summary1)

                st.subheader("GPT 摘要：簡報 2")
                st.markdown(summary2)

                # 額外：比對差異重點（可選）
                with st.spinner("GPT 正在比對兩份摘要差異…"):
                    diff_prompt = (
                        "請以繁體中文比較兩份摘要之差異：\n"
                        "1) 共同主題（要點列出）\n"
                        "2) 各自獨特重點（左/右各列 3-6 點）\n"
                        "3) 建議：如果要合併成一份簡報，建議的結構綱要（3-5 節）\n\n"
                        f"[摘要A]\n{summary1}\n\n[摘要B]\n{summary2}"
                    )
                    diff = client.responses.create(model=model, input=diff_prompt).output_text
                st.subheader("🔍 摘要差異與合併建議")
                st.markdown(diff)

            except Exception as e:
                st.error(f"呼叫 OpenAI 失敗：{e}")
    else:
        st.info("若要產生 GPT 摘要，請先輸入 OpenAI API Key。")
else:
    st.info("請先上傳兩份 PPTX 檔案")
