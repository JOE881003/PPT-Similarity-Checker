import os
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

def extract_text_from_ppt(ppt_path: str) -> str:
    """提取 PPT 中的全部文字"""
    prs = Presentation(ppt_path)
    text_list = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_list.append(shape.text)
    return " ".join(text_list)

def calculate_similarity_main_vs_others(main_ppt_path: str, ppt_paths: list[str]) -> dict:
    """
    計算主 PPT 與一組 PPT 之間的相似度

    Parameters:
        main_ppt_path (str): 主 PPT 檔案路徑
        ppt_paths (list[str]): 要比較的 PPT 路徑清單

    Returns:
        dict: 每個被比較 PPT 對應的相似度分數，例如 {"fileA.pptx": 0.73, ...}
    """
    text_main = extract_text_from_ppt(main_ppt_path)
    texts = [text_main] + [extract_text_from_ppt(p) for p in ppt_paths]
    
    vectorizer = TfidfVectorizer().fit(texts)
    vecs = vectorizer.transform(texts)

    result = {}
    for i, path in enumerate(ppt_paths):
        score = cosine_similarity(vecs[0], vecs[i+1])[0][0]
        result[path] = round(score, 4)
    
    return result
main_path = "example_1.pptx"
others = ["example_2.pptx"]

result = calculate_similarity_main_vs_others(main_path, others)
for f, score in result.items():
    print(f" {main_path} ➝ {f} 相似度: {score}")
